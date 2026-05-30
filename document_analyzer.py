import base64
import csv
import json
import logging
import mimetypes
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
from xml.etree import ElementTree

import replies

logger = logging.getLogger("docs")


@dataclass
class FileContext:
    filename: str
    path: str
    mime_type: str
    content: str
    error: str = ""
    estimated_tokens: int = 0
    token_limit: int = 0
    over_limit: bool = False
    base64_data: str = ""
    source_url: str = ""


class DocumentAnalyzer:
    def __init__(self, max_file_tokens: int = 850000, multimodal: Optional[dict] = None):
        self.max_file_tokens = max_file_tokens
        self.multimodal = multimodal or {}

    def analyze(self, path: Path, mime_type: str = "") -> FileContext:
        mime_type = mime_type or mimetypes.guess_type(path.name)[0] or ""
        if self._is_multimodal(path, mime_type):
            return self._multimodal_context(path, mime_type)

        suffix = path.suffix.lower()
        try:
            if suffix in (".txt", ".md", ".log", ".py", ".js", ".ts", ".html", ".css", ".xml", ".yaml", ".yml"):
                content = self._read_text(path)
            elif suffix == ".json":
                content = self._read_json(path)
            elif suffix == ".csv":
                content = self._read_csv(path)
            elif suffix == ".pdf":
                content = self._read_pdf(path)
            elif suffix == ".docx":
                content = self._read_docx(path)
            elif suffix in (".xlsx", ".xlsm"):
                content = self._read_xlsx(path)
            elif suffix == ".pptx":
                content = self._read_pptx(path)
            else:
                content = ""
                return FileContext(path.name, str(path), mime_type, content, replies.UNSUPPORTED_FORMAT)
            return self._build_context(path, mime_type, content)
        except Exception as exc:
            logger.exception("failed to analyze file: %s", path)
            return FileContext(path.name, str(path), mime_type, "", replies.PARSE_FAILED)

    def _read_text(self, path: Path) -> str:
        for encoding in ("utf-8", "utf-8-sig", "gb18030"):
            try:
                return path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        return path.read_text(errors="replace")

    def _read_json(self, path: Path) -> str:
        data = json.loads(self._read_text(path))
        return json.dumps(data, ensure_ascii=False, indent=2)

    def _read_csv(self, path: Path) -> str:
        import itertools
        with open(path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            preview = list(itertools.islice(reader, 80))
        return "\n".join("\t".join(cell for cell in row) for row in preview)

    def _read_pdf(self, path: Path) -> str:
        try:
            from pypdf import PdfReader
        except ImportError:
            return replies.DEPENDENCY_MISSING_PDF

        reader = PdfReader(str(path))
        pages = []
        for index, page in enumerate(reader.pages[:20], start=1):
            text = page.extract_text() or ""
            pages.append(f"[第 {index} 页]\n{text}")
        return "\n\n".join(pages)

    def _read_docx(self, path: Path) -> str:
        try:
            import docx
        except ImportError:
            return replies.DEPENDENCY_MISSING_DOCX

        document = docx.Document(str(path))
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)

    def _read_xlsx(self, path: Path) -> str:
        try:
            from openpyxl import load_workbook
        except ImportError:
            return replies.DEPENDENCY_MISSING_XLSX

        workbook = load_workbook(str(path), read_only=True, data_only=True)
        chunks = []
        for sheet in workbook.worksheets[:5]:
            chunks.append(f"[工作表：{sheet.title}]")
            for row in sheet.iter_rows(max_row=80, values_only=True):
                values = ["" if cell is None else str(cell) for cell in row[:20]]
                if any(values):
                    chunks.append("\t".join(values))
        workbook.close()
        return "\n".join(chunks)

    def _read_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return self._read_pptx_zip(path)

        presentation = Presentation(str(path))
        chunks = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [cell.text.strip() for cell in row.cells]
                        if any(values):
                            texts.append(" | ".join(values))
            if texts:
                chunks.append(f"[第 {index} 页]\n" + "\n".join(texts))
        return "\n\n".join(chunks)

    def _read_pptx_zip(self, path: Path) -> str:
        chunks = []
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        with zipfile.ZipFile(path) as archive:
            slide_names = sorted(
                [name for name in archive.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)],
                key=lambda name: int(re.search(r"slide(\d+)\.xml$", name).group(1)),
            )
            for index, slide_name in enumerate(slide_names, start=1):
                root = ElementTree.fromstring(archive.read(slide_name))
                texts = []
                for node in root.findall(".//a:t", ns):
                    if node.text and node.text.strip():
                        texts.append(node.text.strip())
                if texts:
                    chunks.append(f"[第 {index} 页]\n" + "\n".join(texts))
        if not chunks:
            return "PPTX 中未提取到可读文本。"
        return "\n\n".join(chunks)

    def _build_context(self, path: Path, mime_type: str, content: str) -> FileContext:
        content = content.strip()
        estimated_tokens = self.estimate_tokens(content)
        if estimated_tokens > self.max_file_tokens:
            return FileContext(
                filename=path.name,
                path=str(path),
                mime_type=mime_type,
                content=content,
                error=replies.FILE_TOO_LONG,
                estimated_tokens=estimated_tokens,
                token_limit=self.max_file_tokens,
                over_limit=True,
            )
        return FileContext(
            filename=path.name,
            path=str(path),
            mime_type=mime_type,
            content=content,
            estimated_tokens=estimated_tokens,
            token_limit=self.max_file_tokens,
        )

    def estimate_tokens(self, text: str) -> int:
        if not text:
            return 0
        cjk = 0
        non_cjk = 0
        for char in text:
            code = ord(char)
            if (
                0x4E00 <= code <= 0x9FFF
                or 0x3400 <= code <= 0x4DBF
                or 0x3040 <= code <= 0x30FF
                or 0xAC00 <= code <= 0xD7AF
            ):
                cjk += 1
            elif not char.isspace():
                non_cjk += 1
        non_cjk_tokens = (non_cjk + 3) // 4 if non_cjk else 0
        return cjk + non_cjk_tokens

    def _is_multimodal(self, path: Path, mime_type: str) -> bool:
        suffix = path.suffix.lower()
        return (
            mime_type.startswith("image/")
            or mime_type.startswith("audio/")
            or mime_type.startswith("video/")
            or suffix in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".mp3", ".wav", ".m4a", ".mp4", ".mov", ".avi", ".mkv")
        )

    def _get_multimodal_status(self) -> str:
        """获取当前多模态功能的状态，用于友好提示"""
        status_lines = []
        # 基础功能
        status_lines.append("  ✅ 文字对话")
        status_lines.append("  ✅ 文档解析（PDF/Word/Excel/TXT）")
        
        # 多模态功能状态
        enabled = self.multimodal.get("enabled", False)
        if not enabled:
            status_lines.append("  ❌ 图片识别（功能未开启）")
            status_lines.append("  ❌ 音频识别（功能未开启）")
            status_lines.append("  ❌ 视频识别（功能未开启）")
        else:
            vision_model = self.multimodal.get("vision_model", "")
            audio_model = self.multimodal.get("audio_model", "")
            video_model = self.multimodal.get("video_model", "") or self.multimodal.get("video_url_model", "")
            
            if vision_model:
                status_lines.append(f"  ✅ 图片识别（{vision_model}）")
            else:
                status_lines.append("  ❌ 图片识别（未配置）")
                
            if audio_model:
                status_lines.append(f"  ✅ 音频识别（{audio_model}）")
            else:
                status_lines.append("  ❌ 音频识别（未配置）")
                
            if video_model:
                status_lines.append(f"  ✅ 视频识别（{video_model}）")
            else:
                status_lines.append("  ❌ 视频识别（未配置）")
                
        return "\n".join(status_lines)

    def _multimodal_context(self, path: Path, mime_type: str) -> FileContext:
        enabled = self.multimodal.get("enabled", False)
        if not enabled:
            status = self._get_multimodal_status()
            return FileContext(
                filename=path.name,
                path=str(path),
                mime_type=mime_type,
                content="",
                error=f"目前多模态功能还没有开启哦～\n\n目前支持的功能：\n{status}\n\n如果需要识别图片/音视频，可以在 config.json 的 multimodal 字段中配置：\n  - 开启功能：设置 enabled 为 true\n  - 图片识别：配置 vision_model（例如 qwen3.6-plus）\n  - 音频识别：配置 audio_model（例如 qwen3.5-omni-plus）\n  - 视频识别：配置 video_model（例如 qwen3.5-omni-plus）",
            )

        is_image = mime_type.startswith("image/") or path.suffix.lower() in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp")
        is_video = mime_type.startswith("video/") or path.suffix.lower() in (".mp4", ".mov", ".avi", ".mkv")
        is_audio = mime_type.startswith("audio/") or path.suffix.lower() in (".mp3", ".wav", ".m4a")

        if is_image:
            max_mb = self.multimodal.get("max_image_mb", 20)
            model = self.multimodal.get("vision_model", "")
        elif is_video:
            max_mb = self.multimodal.get("max_video_mb", 300)
            model = self.multimodal.get("video_model", "") or self.multimodal.get("video_url_model", "")
        elif is_audio:
            max_mb = self.multimodal.get("max_audio_mb", 100)
            model = self.multimodal.get("audio_model", "")
        else:
            max_mb = 20
            model = ""

        file_size_mb = path.stat().st_size / (1024 * 1024)
        if file_size_mb > max_mb:
            return FileContext(
                filename=path.name,
                path=str(path),
                mime_type=mime_type,
                content="",
                error=replies.FILE_TOO_LARGE.format(size=file_size_mb, limit=max_mb),
            )

        if not model:
            if is_image:
                # 智能构建提示，根据当前配置显示哪些功能可用
                status = self._get_multimodal_status()
                error = f"图片识别功能还没配置哦～\n\n目前支持的功能：\n{status}\n\n如果需要识别图片，可以在 config.json 的 multimodal 字段中配置 vision_model（例如 qwen3.6-plus），或者你可以先用文字描述给我～"
            elif is_video:
                status = self._get_multimodal_status()
                error = f"视频识别功能还没配置哦～\n\n目前支持的功能：\n{status}\n\n如果需要识别视频，可以在 config.json 的 multimodal 字段中配置 video_model（例如 qwen3.5-omni-plus），或者你可以发视频链接、截几张关键画面的图发给我～"
            elif is_audio:
                status = self._get_multimodal_status()
                error = f"音频识别功能还没配置哦～\n\n目前支持的功能：\n{status}\n\n如果需要识别音频，可以在 config.json 的 multimodal 字段中配置 audio_model（例如 qwen3.5-omni-plus），或者你可以把想说的打字发给我～"
            else:
                error = replies.UNSUPPORTED_MULTIMODAL
            return FileContext(
                filename=path.name,
                path=str(path),
                mime_type=mime_type,
                content="",
                error=error,
            )

        if is_image:
            try:
                raw = path.read_bytes()
                b64 = base64.b64encode(raw).decode("ascii")
                actual_mime = mime_type or mimetypes.guess_type(path.name)[0] or "image/jpeg"
                return FileContext(
                    filename=path.name,
                    path=str(path),
                    mime_type=actual_mime,
                    content="",
                    base64_data=b64,
                )
            except Exception as exc:
                logger.exception("failed to encode image: %s", path)
                return FileContext(
                    filename=path.name,
                    path=str(path),
                    mime_type=mime_type,
                    content="",
                    error=replies.IMAGE_ENCODE_FAILED,
                )

        if is_audio:
            try:
                raw = path.read_bytes()
                b64 = base64.b64encode(raw).decode("ascii")
                actual_mime = mime_type or mimetypes.guess_type(path.name)[0] or "audio/wav"
                return FileContext(
                    filename=path.name,
                    path=str(path),
                    mime_type=actual_mime,
                    content="",
                    base64_data=b64,
                )
            except Exception as exc:
                logger.exception("failed to encode audio: %s", path)
                return FileContext(
                    filename=path.name,
                    path=str(path),
                    mime_type=mime_type,
                    content="",
                    error=replies.AUDIO_ENCODE_FAILED,
                )

        if is_video:
            try:
                raw = path.read_bytes()
                b64 = base64.b64encode(raw).decode("ascii")
                actual_mime = mime_type or mimetypes.guess_type(path.name)[0] or "video/mp4"
                return FileContext(
                    filename=path.name,
                    path=str(path),
                    mime_type=actual_mime,
                    content="",
                    base64_data=b64,
                )
            except Exception as exc:
                logger.exception("failed to encode video: %s", path)
                return FileContext(
                    filename=path.name,
                    path=str(path),
                    mime_type=mime_type,
                    content="",
                    error=replies.VIDEO_ENCODE_FAILED,
                )

        return FileContext(
            filename=path.name,
            path=str(path),
            mime_type=mime_type,
            content="",
            error=replies.UNSUPPORTED_MULTIMODAL,
        )
